#!/usr/bin/env python
# coding: utf-8

# In[1]:


import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns

# Generar datos para una distribución simétrica
np.random.seed(1)
data_symmetric = np.random.normal(loc=0, scale=1, size=1000)

# Generar datos para una distribución asimétrica hacia la derecha
np.random.seed(2)
data_skew_right = np.random.exponential(scale=2, size=1000)

# Generar datos para una distribución asimétrica hacia la izquierda
np.random.seed(3)
data_skew_left = -np.random.exponential(scale=2, size=1000)

# Crear figura y subplots
fig, axs = plt.subplots(2, 3, figsize=(15, 8))

# Graficar histograma y box plot para la distribución simétrica
axs[0, 0].hist(data_symmetric, bins=20, color='blue', edgecolor='black')
axs[0, 0].set_title('Distribución Simétrica')
axs[0, 0].set_xlabel('Valores')
axs[0, 0].set_ylabel('Frecuencia')
sns.boxplot(data=data_symmetric, ax=axs[0, 1], color='blue')
axs[0, 1].set_title('Box Plot - Simétrico')

# Graficar histograma y box plot para la distribución asimétrica hacia la derecha
axs[0, 2].hist(data_skew_right, bins=20, color='green', edgecolor='black')
axs[0, 2].set_title('Distribución Asimétrica (Derecha)')
axs[0, 2].set_xlabel('Valores')
axs[0, 2].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_right, ax=axs[1, 0], color='green')
axs[1, 0].set_title('Box Plot - Asimétrico (Derecha)')

# Graficar histograma y box plot para la distribución asimétrica hacia la izquierda
axs[1, 1].hist(data_skew_left, bins=20, color='red', edgecolor='black')
axs[1, 1].set_title('Distribución Asimétrica (Izquierda)')
axs[1, 1].set_xlabel('Valores')
axs[1, 1].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_left, ax=axs[1, 2], color='red')
axs[1, 2].set_title('Box Plot - Asimétrico (Izquierda)')

# Ajustar los espacios entre subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()


# In[2]:


import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns

# Generar datos para una distribución simétrica
np.random.seed(1)
data_symmetric = np.random.normal(loc=0, scale=1, size=1000)

# Generar datos para una distribución asimétrica hacia la derecha
np.random.seed(2)
data_skew_right = np.random.exponential(scale=2, size=1000)

# Generar datos para una distribución asimétrica hacia la izquierda
np.random.seed(3)
data_skew_left = -np.random.exponential(scale=2, size=1000)

# Calcular coeficiente de asimetría de Pearson
pearson_skewness_symmetric = 3 * (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
pearson_skewness_skew_right = 3 * (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
pearson_skewness_skew_left = 3 * (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Calcular coeficiente de asimetría de Fisher
fisher_skewness_symmetric = (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
fisher_skewness_skew_right = (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
fisher_skewness_skew_left = (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Crear figura y subplots
fig, axs = plt.subplots(2, 3, figsize=(15, 8))

# Graficar histograma y box plot para la distribución simétrica
axs[0, 0].hist(data_symmetric, bins=20, color='blue', edgecolor='black')
axs[0, 0].set_title('Distribución Simétrica')
axs[0, 0].set_xlabel('Valores')
axs[0, 0].set_ylabel('Frecuencia')
sns.boxplot(data=data_symmetric, ax=axs[0, 1], color='blue')
axs[0, 1].set_title('Box Plot - Simétrico')

# Graficar histograma y box plot para la distribución asimétrica hacia la derecha
axs[0, 2].hist(data_skew_right, bins=20, color='green', edgecolor='black')
axs[0, 2].set_title('Distribución Asimétrica (Derecha)')
axs[0, 2].set_xlabel('Valores')
axs[0, 2].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_right, ax=axs[1, 0], color='green')
axs[1, 0].set_title('Box Plot - Asimétrico (Derecha)')

# Graficar histograma y box plot para la distribución asimétrica hacia la izquierda
axs[1, 1].hist(data_skew_left, bins=20, color='red', edgecolor='black')
axs[1, 1].set_title('Distribución Asimétrica (Izquierda)')
axs[1, 1].set_xlabel('Valores')
axs[1, 1].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_left, ax=axs[1, 2], color='red')
axs[1, 2].set_title('Box Plot - Asimétrico (Izquierda)')

# Mostrar coeficientes de asimetría
axs[0, 0].text(0.05, 0.95, f"Pearson: {pearson_skewness_symmetric:.2f}", transform=axs[0, 0].transAxes, verticalalignment='top', color='black')
axs[0, 1].text(0.05, 0.95, f"Pearson: {pearson_skewness_symmetric:.2f}", transform=axs[0, 1].transAxes, verticalalignment='top', color='black')
axs[0, 2].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_right:.2f}", transform=axs[0, 2].transAxes, verticalalignment='top', color='black')
axs[1, 0].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_right:.2f}", transform=axs[1, 0].transAxes, verticalalignment='top', color='black')
axs[1, 1].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_left:.2f}", transform=axs[1, 1].transAxes, verticalalignment='top', color='black')
axs[1, 2].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_left:.2f}", transform=axs[1, 2].transAxes, verticalalignment='top', color='black')

# Ajustar los espacios entre subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()


# In[11]:


import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns

# Generar datos para una distribución simétrica
np.random.seed(1)
data_symmetric = np.random.normal(loc=0, scale=1, size=1000)

# Generar datos para una distribución asimétrica hacia la derecha
np.random.seed(2)
data_skew_right = np.random.exponential(scale=2, size=1000)

# Generar datos para una distribución asimétrica hacia la izquierda
np.random.seed(3)
data_skew_left = -np.random.exponential(scale=2, size=1000)

# Calcular coeficiente de asimetría de Pearson
pearson_skewness_symmetric = 3 * (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
pearson_skewness_skew_right = 3 * (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
pearson_skewness_skew_left = 3 * (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Calcular coeficiente de asimetría de Fisher
fisher_skewness_symmetric = (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
fisher_skewness_skew_right = (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
fisher_skewness_skew_left = (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Crear figura y subplots
fig, axs = plt.subplots(2, 3, figsize=(15, 8))

# Graficar histograma y box plot para la distribución simétrica
axs[0, 0].hist(data_symmetric, bins=20, color='blue', edgecolor='black')
axs[0, 0].set_title('Distribución Simétrica')
axs[0, 0].set_xlabel('Valores')
axs[0, 0].set_ylabel('Frecuencia')
sns.boxplot(data=data_symmetric, ax=axs[0, 1], color='blue')
axs[0, 1].set_title('Box Plot - Simétrico')

# Graficar histograma y box plot para la distribución asimétrica hacia la derecha
axs[0, 2].hist(data_skew_right, bins=20, color='green', edgecolor='black')
axs[0, 2].set_title('Distribución Asimétrica (Derecha)')
axs[0, 2].set_xlabel('Valores')
axs[0, 2].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_right, ax=axs[1, 0], color='green')
axs[1, 0].set_title('Box Plot - Asimétrico (Derecha)')

# Graficar histograma y box plot para la distribución asimétrica hacia la izquierda
axs[1, 1].hist(data_skew_left, bins=20, color='red', edgecolor='black')
axs[1, 1].set_title('Distribución Asimétrica (Izquierda)')
axs[1, 1].set_xlabel('Valores')
axs[1, 1].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_left, ax=axs[1, 2], color='red')
axs[1, 2].set_title('Box Plot - Asimétrico (Izquierda)')

# Mostrar coeficientes de asimetría
axs[0, 0].text(0.05, 0.95, f"Pearson: {pearson_skewness_symmetric:.2f}\nFisher: {fisher_skewness_symmetric:.2f}", transform=axs[0, 0].transAxes, verticalalignment='top', color='black')
axs[0, 1].text(0.05, 0.95, f"Pearson: {pearson_skewness_symmetric:.2f}\nFisher: {fisher_skewness_symmetric:.2f}", transform=axs[0, 1].transAxes, verticalalignment='top', color='black')
axs[0, 2].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_right:.2f}\nFisher: {fisher_skewness_skew_right:.2f}", transform=axs[0, 2].transAxes, verticalalignment='top', color='black')
axs[1, 0].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_right:.2f}\nFisher: {fisher_skewness_skew_right:.2f}", transform=axs[1, 0].transAxes, verticalalignment='top', color='black')
axs[1, 1].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_left:.2f}\nFisher: {fisher_skewness_skew_left:.2f}", transform=axs[1, 1].transAxes, verticalalignment='top', color='black')
axs[1, 2].text(0.05, 0.95, f"Pearson: {pearson_skewness_skew_left:.2f}\nFisher: {fisher_skewness_skew_left:.2f}", transform=axs[1, 2].transAxes, verticalalignment='top', color='black')

# Ajustar los espacios entre subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()


# In[12]:


import numpy as np

# Generar datos para una distribución simétrica
np.random.seed(1)
data_symmetric = np.random.normal(loc=0, scale=1, size=1000)

# Generar datos para una distribución asimétrica hacia la derecha
np.random.seed(2)
data_skew_right = np.random.exponential(scale=2, size=1000)

# Generar datos para una distribución asimétrica hacia la izquierda
np.random.seed(3)
data_skew_left = -np.random.exponential(scale=2, size=1000)

# Calcular coeficiente de asimetría de Pearson
pearson_skewness_symmetric = 3 * (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
pearson_skewness_skew_right = 3 * (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
pearson_skewness_skew_left = 3 * (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Calcular coeficiente de asimetría de Fisher
fisher_skewness_symmetric = (np.mean(data_symmetric) - np.median(data_symmetric)) / np.std(data_symmetric)
fisher_skewness_skew_right = (np.mean(data_skew_right) - np.median(data_skew_right)) / np.std(data_skew_right)
fisher_skewness_skew_left = (np.mean(data_skew_left) - np.median(data_skew_left)) / np.std(data_skew_left)

# Imprimir los resultados
print("Coeficientes de asimetría:")
print(f"Distribución simétrica - Pearson: {pearson_skewness_symmetric:.2f}, Fisher: {fisher_skewness_symmetric:.2f}")
print(f"Distribución asimétrica hacia la derecha - Pearson: {pearson_skewness_skew_right:.2f}, Fisher: {fisher_skewness_skew_right:.2f}")
print(f"Distribución asimétrica hacia la izquierda - Pearson: {pearson_skewness_skew_left:.2f}, Fisher: {fisher_skewness_skew_left:.2f}")


# In[6]:


import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import kurtosis

# Generar datos para una distribución simétrica
np.random.seed(1)
data_symmetric = np.random.normal(loc=0, scale=1, size=1000)

# Generar datos para una distribución asimétrica hacia la derecha
np.random.seed(2)
data_skew_right = np.random.exponential(scale=2, size=1000)

# Generar datos para una distribución asimétrica hacia la izquierda
np.random.seed(3)
data_skew_left = -np.random.exponential(scale=2, size=1000)

# Calcular coeficientes de curtosis
kurtosis_symmetric = np.round(kurtosis(data_symmetric), 2)
kurtosis_skew_right = np.round(kurtosis(data_skew_right), 2)
kurtosis_skew_left = np.round(kurtosis(data_skew_left), 2)

# Crear figura y subplots
fig, axs = plt.subplots(2, 3, figsize=(15, 8))

# Graficar histograma y box plot para la distribución simétrica
axs[0, 0].hist(data_symmetric, bins=20, color='blue', edgecolor='black')
axs[0, 0].set_title('Distribución Simétrica')
axs[0, 0].set_xlabel('Valores')
axs[0, 0].set_ylabel('Frecuencia')
sns.boxplot(data=data_symmetric, ax=axs[0, 1], color='blue')
axs[0, 1].set_title('Box Plot - Simétrico')

# Graficar histograma y box plot para la distribución asimétrica hacia la derecha
axs[0, 2].hist(data_skew_right, bins=20, color='green', edgecolor='black')
axs[0, 2].set_title('Distribución Asimétrica (Derecha)')
axs[0, 2].set_xlabel('Valores')
axs[0, 2].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_right, ax=axs[1, 0], color='green')
axs[1, 0].set_title('Box Plot - Asimétrico (Derecha)')

# Graficar histograma y box plot para la distribución asimétrica hacia la izquierda
axs[1, 1].hist(data_skew_left, bins=20, color='red', edgecolor='black')
axs[1, 1].set_title('Distribución Asimétrica (Izquierda)')
axs[1, 1].set_xlabel('Valores')
axs[1, 1].set_ylabel('Frecuencia')
sns.boxplot(data=data_skew_left, ax=axs[1, 2], color='red')
axs[1, 2].set_title('Box Plot - Asimétrico (Izquierda)')

# Agregar texto con coeficiente de curtosis en cada gráfico
axs[0, 0].text(0.05, 0.9, f'Curtosis: {kurtosis_symmetric}', transform=axs[0, 0].transAxes, fontsize=10)
axs[0, 2].text(0.05, 0.9, f'Curtosis: {kurtosis_skew_right}', transform=axs[0, 2].transAxes, fontsize=10)
axs[1, 1].text(0.05, 0.9, f'Curtosis: {kurtosis_skew_left}', transform=axs[1, 1].transAxes, fontsize=10)

# Ajustar los espacios entre subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()


# In[7]:


import pandas as pd
import seaborn as sns
tips = sns.load_dataset('tips')
tips


# In[10]:


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Crear figura y subplots
fig, axs = plt.subplots(1, 2, figsize=(12, 6))

# Graficar histograma del total de la cuenta (total_bill)
axs[0].hist(tips['total_bill'], bins=20, color='blue', edgecolor='black')
axs[0].set_title('Distribución del Total de la Cuenta')
axs[0].set_xlabel('Total de la Cuenta')
axs[0].set_ylabel('Frecuencia')

# Graficar histograma de la propina (tip)
axs[1].hist(tips['tip'], bins=20, color='green', edgecolor='black')
axs[1].set_title('Distribución de la Propina')
axs[1].set_xlabel('Propina')
axs[1].set_ylabel('Frecuencia')

# Agregar texto con simetría/asimetría y curtosis en cada gráfico
axs[0].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[0].text(0.05, 0.8, f'Curtosis: {kurtosis_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[1].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_tip}', transform=axs[1].transAxes, fontsize=10)
axs[1].text(0.05, 0.8, f'Curtosis: {kurtosis_tip}', transform=axs[1].transAxes, fontsize=10)

# Crear subplot adicional para el boxplot
fig, ax = plt.subplots(figsize=(6, 6))

# Graficar boxplot del total de la cuenta (total_bill)
sns.boxplot(data=tips, x='total_bill', ax=ax, color='blue')
ax.set_title('Box Plot - Total de la Cuenta')
ax.set_xlabel('Total de la Cuenta')

# Mostrar los gráficos
plt.tight_layout()
plt.show()


# In[8]:


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Crear figura y subplots
fig, axs = plt.subplots(1, 2, figsize=(12, 6))

# Graficar histograma del total de la cuenta (total_bill)
axs[0].hist(tips['total_bill'], bins=20, color='blue', edgecolor='black')
axs[0].set_title('Distribución del Total de la Cuenta')
axs[0].set_xlabel('Total de la Cuenta')
axs[0].set_ylabel('Frecuencia')

# Graficar histograma de la propina (tip)
axs[1].hist(tips['tip'], bins=20, color='green', edgecolor='black')
axs[1].set_title('Distribución de la Propina')
axs[1].set_xlabel('Propina')
axs[1].set_ylabel('Frecuencia')

# Agregar texto con simetría/asimetría y curtosis en cada gráfico
axs[0].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[0].text(0.05, 0.8, f'Curtosis: {kurtosis_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[1].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_tip}', transform=axs[1].transAxes, fontsize=10)
axs[1].text(0.05, 0.8, f'Curtosis: {kurtosis_tip}', transform=axs[1].transAxes, fontsize=10)

# Ajustar los espacios entre subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()



# In[9]:


import pandas as pd
import numpy as np
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Imprimir los resultados
print("Simetría/Asimetría:")
print(f"Total de la Cuenta: {skewness_total_bill}")
print(f"Propina: {skewness_tip}")
print()
print("Curtosis:")
print(f"Total de la Cuenta: {kurtosis_total_bill}")
print(f"Propina: {kurtosis_tip}")


# In[13]:


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Crear figura y subplots
fig, axs = plt.subplots(1, 2, figsize=(12, 6))

# Graficar histograma del total de la cuenta (total_bill)
axs[0].hist(tips['total_bill'], bins=20, color='blue', edgecolor='black')
axs[0].set_title('Distribución del Total de la Cuenta')
axs[0].set_xlabel('Total de la Cuenta')
axs[0].set_ylabel('Frecuencia')

# Graficar histograma de la propina (tip)
axs[1].hist(tips['tip'], bins=20, color='green', edgecolor='black')
axs[1].set_title('Distribución de la Propina')
axs[1].set_xlabel('Propina')
axs[1].set_ylabel('Frecuencia')

# Agregar texto con simetría/asimetría y curtosis en cada gráfico
axs[0].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[0].text(0.05, 0.8, f'Curtosis: {kurtosis_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[1].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_tip}', transform=axs[1].transAxes, fontsize=10)
axs[1].text(0.05, 0.8, f'Curtosis: {kurtosis_tip}', transform=axs[1].transAxes, fontsize=10)

# Crear subplot adicional para el boxplot
fig, ax = plt.subplots(figsize=(6, 6))

# Graficar boxplot del total de la cuenta (total_bill)
sns.boxplot(data=tips, x='total_bill', ax=ax, color='blue')
ax.set_title('Box Plot - Total de la Cuenta')
ax.set_xlabel('Total de la Cuenta')

# Ruta de destino para guardar los gráficos
ruta_destino = r'C:\Users\Sergio\Documents\SGdataconsulting\Youtube\Estadística'

# Guardar los gráficos en la ruta especificada
plt.savefig(ruta_destino + '/histogramas.png')
plt.savefig(ruta_destino + '/boxplot.png')

# Mostrar los gráficos
plt.tight_layout()
plt.show()


# In[14]:


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Crear figura y subplots
fig, axs = plt.subplots(1, 2, figsize=(12, 6))

# Graficar histograma del total de la cuenta (total_bill)
axs[0].hist(tips['total_bill'], bins=20, color='blue', edgecolor='black')
axs[0].set_title('Distribución del Total de la Cuenta')
axs[0].set_xlabel('Total de la Cuenta')
axs[0].set_ylabel('Frecuencia')

# Graficar boxplot del total de la cuenta (total_bill)
sns.boxplot(data=tips, x='total_bill', ax=axs[1], color='blue')
axs[1].set_title('Box Plot - Total de la Cuenta')
axs[1].set_xlabel('Total de la Cuenta')

# Agregar texto con simetría/asimetría y curtosis en la figura
fig.text(0.05, 0.9, f'Simetría/Asimetría: {skewness_total_bill}', fontsize=10)
fig.text(0.05, 0.85, f'Curtosis: {kurtosis_total_bill}', fontsize=10)

# Ruta de destino para guardar la figura combinada
ruta_destino = r'C:\Users\Sergio\Documents\SGdataconsulting\Youtube\Estadística'

# Guardar la figura combinada como JPEG
plt.savefig(ruta_destino + '/histograma_boxplot.jpeg', format='jpeg')

# Mostrar la figura
plt.tight_layout()
plt.show()


# In[15]:


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import skew, kurtosis

# Cargar el dataset 'tips'
tips = sns.load_dataset('tips')

# Calcular la simetría o asimetría
skewness_total_bill = np.round(skew(tips['total_bill']), 2)
skewness_tip = np.round(skew(tips['tip']), 2)

# Calcular la curtosis
kurtosis_total_bill = np.round(kurtosis(tips['total_bill']), 2)
kurtosis_tip = np.round(kurtosis(tips['tip']), 2)

# Crear figura y subplots para los histogramas
fig, axs = plt.subplots(2, 1, figsize=(8, 10))

# Graficar histograma del total de la cuenta (total_bill)
axs[0].hist(tips['total_bill'], bins=20, color='blue', edgecolor='black')
axs[0].set_title('Distribución del Total de la Cuenta')
axs[0].set_xlabel('Total de la Cuenta')
axs[0].set_ylabel('Frecuencia')
axs[0].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_total_bill}', transform=axs[0].transAxes, fontsize=10)
axs[0].text(0.05, 0.8, f'Curtosis: {kurtosis_total_bill}', transform=axs[0].transAxes, fontsize=10)

# Graficar histograma de la propina (tip)
axs[1].hist(tips['tip'], bins=20, color='green', edgecolor='black')
axs[1].set_title('Distribución de la Propina')
axs[1].set_xlabel('Propina')
axs[1].set_ylabel('Frecuencia')
axs[1].text(0.05, 0.9, f'Simetría/Asimetría: {skewness_tip}', transform=axs[1].transAxes, fontsize=10)
axs[1].text(0.05, 0.8, f'Curtosis: {kurtosis_tip}', transform=axs[1].transAxes, fontsize=10)

# Crear figura y subplots para los boxplots
fig, axs = plt.subplots(2, 1, figsize=(8, 10))

# Graficar boxplot del total de la cuenta (total_bill)
sns.boxplot(data=tips, y='total_bill', ax=axs[0], color='blue')
axs[0].set_title('Box Plot - Total de la Cuenta')
axs[0].set_ylabel('Total de la Cuenta')

# Graficar boxplot de la propina (tip)
sns.boxplot(data=tips, y='tip', ax=axs[1], color='green')
axs[1].set_title('Box Plot - Propina')
axs[1].set_ylabel('Propina')

# Ajustar los espacios entre los subplots
plt.tight_layout()

# Mostrar los gráficos
plt.show()


# In[18]:


from pptx import Presentation

# Crear una presentación
presentation = Presentation()

# Agregar diapositiva de título y contenido
slide_layout = presentation.slide_layouts[1]  # Diseño de título y contenido
slide = presentation.slides.add_slide(slide_layout)

# Agregar título
title = slide.shapes.title
title.text = "Agradecimientos"

# Agregar contenido
content = slide.placeholders[1]
content.text = "Quiero agradecer a todos los que hicieron posible este proyecto."

# Agregar diapositiva de título y contenido para el cierre
slide_layout = presentation.slide_layouts[1]  # Diseño de título y contenido
slide = presentation.slides.add_slide(slide_layout)

# Agregar título para el cierre
title = slide.shapes.title
title.text = "Cierre de la presentación"

# Agregar contenido para el cierre
content = slide.placeholders[1]
content.text = "Gracias por su atención. ¿Hay alguna pregunta?"

# Guardar la presentación
ruta_destino = r'C:\Users\Sergio\Documents\SGdataconsulting\Youtube\Estadística\presentacion.pptx'
presentation.save(ruta_destino)



# In[21]:


jupyter nbconvert 'Estadística- Medidas de forma.ipynb' --to slides --post serve



# In[ ]:




